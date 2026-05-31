"""
Genera la presentación PI de Bomb Champions (apoyo oral, texto mínimo).
Ejecutar: .venv\\Scripts\\python.exe PRESENTACION\\generar_presentacion.py
"""

from __future__ import annotations

import re
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
PRESENTACION = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
SALIDA = PRESENTACION / "proyecto intermodular.pptx"
NOTAS = PRESENTACION / "MARCADORES_presentacion.md"
EVOLUCION_DOCX = PRESENTACION / "Capturas de pantallad de evolucion..docx"
EVOLUCION_DIR = PRESENTACION / "ENTREGA" / "imagenes" / "evolucion"
CAPTURAS_DIR = PRESENTACION / "ENTREGA" / "imagenes" / "capturas"

# Estilo (coherente con diapositivas Canva existentes)
COLOR_FONDO = RGBColor(8, 8, 12)
COLOR_MORADO = RGBColor(130, 90, 220)
COLOR_BLANCO = RGBColor(245, 245, 250)
COLOR_GRIS = RGBColor(170, 170, 185)
COLOR_MARCADOR = RGBColor(255, 210, 70)
COLOR_MARCADOR_TEXTO = RGBColor(40, 35, 20)
FUENTE = "Arial"


def cargar_capturas_evolucion() -> list[Path]:
    extraidas = extraer_imagenes_docx(EVOLUCION_DOCX, EVOLUCION_DIR, "evolucion")
    if extraidas:
        return extraidas
    if CAPTURAS_DIR.exists():
        return sorted(CAPTURAS_DIR.glob("evolucion_*.*"))
    return []


def captura_evolucion(capturas: list[Path], indice: int) -> Path | None:
    if 0 <= indice < len(capturas):
        return capturas[indice]
    return None


def extraer_imagenes_docx(docx_path: Path, out_dir: Path, prefix: str) -> list[Path]:
    if not docx_path.exists():
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    guardadas: list[Path] = []
    with zipfile.ZipFile(docx_path) as z:
        rels = {}
        for rel in ET.fromstring(z.read("word/_rels/document.xml.rels")):
            target = rel.get("Target")
            if target and "media/" in target:
                rels[rel.get("Id")] = "word/" + target.replace("../", "")
        xml = z.read("word/document.xml").decode("utf-8")
        order = re.findall(r'r:embed="(rId\d+)"', xml)
        for i, rid in enumerate(order, 1):
            media = rels.get(rid)
            if not media:
                continue
            dest = out_dir / f"{prefix}_{i:02d}{Path(media).suffix}"
            dest.write_bytes(z.read(media))
            guardadas.append(dest)
    return guardadas


def fondo_oscuro(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO


def brillo_morado(slide, izquierda: bool = True):
    forma = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(-1.2 if izquierda else 8.5),
        Inches(-0.8 if izquierda else 4.2),
        Inches(4.5),
        Inches(4.5),
    )
    forma.fill.solid()
    forma.fill.fore_color.rgb = COLOR_MORADO
    forma.fill.transparency = 0.82
    forma.line.fill.background()


def titulo(slide, texto: str, tamano: int = 32):
    caja = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9))
    marco = caja.text_frame
    marco.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = marco.paragraphs[0]
    p.text = texto.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FUENTE
    p.font.size = Pt(tamano)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO


def subtitulo(slide, texto: str, y: float = 1.35, tamano: int = 20):
    caja = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.6))
    p = caja.text_frame.paragraphs[0]
    p.text = texto
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FUENTE
    p.font.size = Pt(tamano)
    p.font.color.rgb = COLOR_GRIS


def keywords(slide, items: list[str], y: float = 2.2):
    texto = "  ·  ".join(items)
    caja = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(1.2))
    p = caja.text_frame.paragraphs[0]
    p.text = texto
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FUENTE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_MORADO


def texto_corto(slide, texto: str, x: float, y: float, w: float, h: float, tamano: int = 16):
    caja = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = caja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.name = FUENTE
    p.font.size = Pt(tamano)
    p.font.color.rgb = COLOR_GRIS


def imagen(slide, ruta: Path, x: float, y: float, w: float):
    if ruta.exists():
        slide.shapes.add_picture(str(ruta), Inches(x), Inches(y), width=Inches(w))
        return True
    return False


def marcador(slide, instruccion: str, x: float, y: float, w: float, h: float):
    forma = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    forma.fill.solid()
    forma.fill.fore_color.rgb = COLOR_MARCADOR
    forma.line.color.rgb = COLOR_MORADO
    forma.line.width = Pt(1.5)
    tf = forma.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "[INSERTAR IMAGEN]\n" + instruccion
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FUENTE
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_MARCADOR_TEXTO


def nueva_diapositiva(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fondo_oscuro(slide)
    brillo_morado(slide, izquierda=True)
    brillo_morado(slide, izquierda=False)
    return slide


def diapositiva_portada(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Bomb Champions", tamano=44)
    subtitulo(slide, "Proyecto Intermodular · SMR", y=1.5, tamano=18)
    subtitulo(slide, "Miguel Marcano · Gabriel Celis", y=2.15, tamano=22)
    subtitulo(slide, "IES El Álamo · 5 junio 2026", y=6.5, tamano=16)
    if ASSETS.joinpath("FONDO BOMB CHAMPIOS.png").exists():
        imagen(slide, ASSETS / "FONDO BOMB CHAMPIOS.png", 7.2, 1.0, 5.5)
    return slide


def diapositiva_equipo(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Equipo")
    keywords(slide, ["Miguel → código · redes", "Gabriel → sprites · audio"], y=1.6)
    texto_corto(
        slide,
        "Python · Pygame · PyCharm          Piskel · BeepBox",
        0.8,
        2.8,
        11.7,
        0.5,
        18,
    )
    for i, (nombre, archivo) in enumerate(
        [("Miguel", "BLUE ABAJO.png"), ("Gabriel", "CuchillasPJ ABAJO.png")]
    ):
        x = 2.2 + i * 5.5
        imagen(slide, ASSETS / archivo, x, 3.5, 1.4)
        caja = slide.shapes.add_textbox(Inches(x), Inches(5.2), Inches(2.5), Inches(0.4))
        p = caja.text_frame.paragraphs[0]
        p.text = nombre
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FUENTE
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO


def diapositiva_juego(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "¿Qué es Bomb Champions?")
    keywords(slide, ["Bomberman · Campeones · LAN"], y=1.5)
    texto_corto(
        slide,
        "Laberinto · bombas · último en pie",
        1.0,
        2.3,
        11.3,
        0.5,
        20,
    )
    imagen(slide, ASSETS / "FONDO BOMB CHAMPIOS.png", 0.8, 3.0, 11.7)


def diapositiva_objetivos(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Objetivos")
    keywords(
        slide,
        ["Videojuego funcional", "POO + Pygame", "Multijugador LAN", "Assets propios"],
        y=2.0,
    )
    texto_corto(slide, "MVP jugable y documentado", 0.8, 3.2, 11.7, 0.5, 18)


def diapositiva_modulos(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Integración intermodular")
    keywords(slide, ["Python → lógica del juego", "Redes → host/cliente LAN"], y=1.8)
    marcador(
        slide,
        "Diagrama simple: 2 PCs en Wi‑Fi → host (TCP) + descubrimiento (UDP).\n"
        "Puedes dibujarlo en Canva o exportar un esquema de la memoria.",
        2.5,
        3.0,
        8.3,
        3.5,
    )


def diapositiva_arquitectura(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Arquitectura del código")
    modulos = (
        "main.py\n"
        "├── configuracion.py\n"
        "├── mapa.py · bomba.py · campeones.py\n"
        "├── especiales.py · hud.py\n"
        "└── red_descubrimiento.py · red_partida.py"
    )
    caja = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.5), Inches(4.5))
    tf = caja.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = modulos
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_BLANCO
    marcador(
        slide,
        "Diagrama de módulos (cajas y flechas) como en la memoria:\n"
        "PRESENTACION/ENTREGA/imagenes/diagramas/arquitectura_modulos.png",
        6.5,
        1.6,
        6.0,
        4.2,
    )


def diapositiva_mapa(prs: Presentation, capturas: list[Path]):
    slide = nueva_diapositiva(prs)
    titulo(slide, "El mapa")
    keywords(slide, ["17 × 21 celdas", "Suelo · pared · caja"], y=1.5)
    partida = captura_evolucion(capturas, 5)
    if partida and partida.exists():
        imagen(slide, partida, 0.7, 2.0, 8.8)
    else:
        marcador(
            slide,
            "Captura de partida en curso (mapa completo con jugadores).",
            0.7,
            2.0,
            8.8,
            4.5,
        )
    for i, (archivo, etiqueta) in enumerate(
        [
            ("Pasto.png", "Suelo"),
            ("ParedHierro.gif", "Pared"),
            ("ParedLadrillos.png", "Caja"),
        ]
    ):
        y = 2.2 + i * 1.5
        imagen(slide, ASSETS / archivo, 10.0, y, 1.2)
        caja = slide.shapes.add_textbox(Inches(11.3), Inches(y + 0.15), Inches(1.5), Inches(0.4))
        p = caja.text_frame.paragraphs[0]
        p.text = etiqueta
        p.font.name = FUENTE
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_GRIS


def diapositiva_gameplay(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Gameplay")
    keywords(slide, ["Bombas · explosión · vidas · invencibilidad"], y=1.5)
    imagen(slide, ASSETS / "Bomba.png", 1.5, 2.4, 2.0)
    imagen(slide, ASSETS / "EXPLOCION.png", 4.2, 2.4, 3.5)
    imagen(slide, ASSETS / "CORAZON.png", 8.5, 2.6, 0.8)
    texto_corto(slide, "3 vidas · HUD por jugador", 8.0, 3.6, 4.0, 0.5, 14)


def diapositiva_campeones(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Campeones")
    campeones = [
        ("BLUE ABAJO.png", "Azul"),
        ("RED ABAJO.png", "Rojo"),
        ("YELLOW ABAJO.png", "Amarillo"),
        ("CuchillasPJ ABAJO.png", "Cuchillas"),
    ]
    for i, (archivo, nombre) in enumerate(campeones):
        x = 0.9 + i * 3.1
        imagen(slide, ASSETS / archivo, x, 2.0, 1.3)
        caja = slide.shapes.add_textbox(Inches(x - 0.3), Inches(4.0), Inches(2.0), Inches(0.4))
        p = caja.text_frame.paragraphs[0]
        p.text = nombre
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FUENTE
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BLANCO
    keywords(slide, ["Cada uno con stats y especial distinto"], y=4.8)
    imagen(slide, ASSETS / "HABILIDAD CUCHILLO DERECHA.png", 9.5, 4.6, 2.5)


def diapositiva_piskel(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Diseño gráfico · Piskel")
    keywords(slide, ["Pixel art · sprites · tiles"], y=1.5)
    texto_corto(slide, "Gabriel Celis", 0.8, 2.1, 11.7, 0.4, 16)
    for i, archivo in enumerate(["Pasto.png", "Bomba.png", "ParedLadrillos.png"]):
        imagen(slide, ASSETS / archivo, 1.0 + i * 3.5, 2.8, 2.5)
    subtitulo(slide, "Fondo menú con apoyo de IA (ChatGPT)", y=5.6, tamano=14)


def diapositiva_audio(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Música y sonido · BeepBox")
    keywords(slide, ["Partida · bomba · menú"], y=1.8)
    marcador(
        slide,
        "Logo de BeepBox (beepbox.co) o captura de la pista del instrumental.\n"
        "Opcional: icono de altavoz + nombre de los 3 archivos .mid del proyecto.",
        2.0,
        2.8,
        9.3,
        3.5,
    )


def diapositiva_dificultad_gif(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Dificultad: animaciones GIF")
    keywords(slide, ["Problema → spritesheet PNG"], y=1.6)
    texto_corto(
        slide,
        "GIF = 1 frame en Pygame · Piskel falla con muchos colores",
        0.8,
        2.3,
        11.5,
        0.8,
        16,
    )
    texto_corto(
        slide,
        "Solución: exportar frames PNG y recortar en código",
        0.8,
        3.0,
        11.5,
        0.5,
        18,
    )
    imagen(slide, ASSETS / "BLUE DERECHA.gif", 1.5, 3.8, 2.0)
    imagen(slide, ASSETS / "BLUE DERECHA.png", 5.0, 3.8, 2.0)


def diapositiva_red(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Multijugador LAN")
    keywords(slide, ["Crear partida · Buscar · Sincronizar"], y=1.6)
    texto_corto(
        slide,
        "UDP descubrimiento · TCP partida · host autoridad",
        0.8,
        2.3,
        11.5,
        0.5,
        16,
    )
    marcador(
        slide,
        "Captura del menú «Crear partida» mostrando la IP del host,\n"
        "o foto de 2 portátiles jugando en la misma red.",
        1.5,
        3.0,
        10.3,
        3.5,
    )


def diapositiva_evolucion(prs: Presentation, capturas: list[Path]):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Evolución del juego")
    # evolucion_03 = MVP · evolucion_05 = selección campeones · evolucion_06 = actual
    slots = [(2, "v1 MVP"), (4, "v2 menú"), (5, "Actual")]
    for i, (idx, etiqueta) in enumerate(slots):
        x = 0.8 + i * 4.2
        ruta = captura_evolucion(capturas, idx)
        if ruta and ruta.exists():
            imagen(slide, ruta, x, 1.5, 3.6)
        else:
            marcador(
                slide,
                f"Captura {etiqueta} del docx de evolución\n"
                "(PRESENTACION/Capturas de pantallad de evolucion..docx)",
                x,
                1.5,
                3.6,
                3.5,
            )
        caja = slide.shapes.add_textbox(Inches(x), Inches(5.3), Inches(3.6), Inches(0.4))
        p = caja.text_frame.paragraphs[0]
        p.text = etiqueta
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FUENTE
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_GRIS


def diapositiva_demo(prs: Presentation, capturas: list[Path]):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Demo y entrega")
    keywords(slide, ["dist/BombChampions.exe · portable"], y=1.6)
    menu = captura_evolucion(capturas, 3)  # evolucion_04: menú con arte
    github = captura_evolucion(capturas, 6)  # evolucion_07: repositorio
    if menu and imagen(slide, menu, 0.8, 2.3, 6.0):
        pass
    else:
        marcador(
            slide,
            "Captura del menú principal en ejecución (demo en vivo en la exposición).",
            0.8,
            2.3,
            6.0,
            3.8,
        )
    if github and imagen(slide, github, 7.2, 2.3, 5.5):
        pass
    else:
        marcador(
            slide,
            "Captura del repositorio GitHub del proyecto.",
            7.2,
            2.3,
            5.5,
            3.8,
        )


def diapositiva_ia(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "IA como apoyo")
    keywords(slide, ["Depurar · organizar · redactar borradores"], y=1.8)
    texto_corto(
        slide,
        "Siempre revisado y probado por el equipo",
        0.8,
        2.8,
        11.5,
        0.5,
        18,
    )


def diapositiva_conclusiones(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "Conclusiones")
    keywords(
        slide,
        ["Base jugable", "Python + redes + diseño", "Mejoras futuras"],
        y=2.0,
    )
    texto_corto(
        slide,
        "Más campeones · IA enemigos · mapas nuevos",
        0.8,
        3.2,
        11.5,
        0.5,
        16,
    )


def diapositiva_preguntas(prs: Presentation):
    slide = nueva_diapositiva(prs)
    titulo(slide, "¿Preguntas?", tamano=48)
    subtitulo(slide, "Gracias", y=3.0, tamano=28)
    imagen(slide, ASSETS / "Bomba.png", 5.8, 4.2, 1.2)


def escribir_notas_marcadores(notas: list[tuple[int, str, str]]):
    lineas = [
        "# Marcadores de la presentación",
        "",
        "Diapositivas con recuadros amarillos: sustituir por la imagen indicada.",
        "Regenerar el PPTX: `python PRESENTACION/generar_presentacion.py`",
        "",
        "| Slide | Título | Qué insertar manualmente |",
        "|-------|--------|--------------------------|",
    ]
    for num, titulo_slide, instruccion in notas:
        lineas.append(f"| {num} | {titulo_slide} | {instruccion} |")
    NOTAS.write_text("\n".join(lineas) + "\n", encoding="utf-8")


def main():
    capturas = cargar_capturas_evolucion()

    backup = PRESENTACION / "proyecto intermodular_backup.pptx"
    if SALIDA.exists() and not backup.exists():
        shutil.copy2(SALIDA, backup)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    marcadores: list[tuple[int, str, str]] = []

    diapositiva_portada(prs)
    diapositiva_equipo(prs)
    diapositiva_juego(prs)
    diapositiva_objetivos(prs)
    diapositiva_modulos(prs)
    marcadores.append((5, "Integración intermodular", "Diagrama LAN (2 PCs + host/cliente)"))
    diapositiva_arquitectura(prs)
    marcadores.append((6, "Arquitectura", "Diagrama arquitectura_modulos.png de la memoria"))
    diapositiva_mapa(prs, capturas)
    diapositiva_gameplay(prs)
    diapositiva_campeones(prs)
    diapositiva_piskel(prs)
    diapositiva_audio(prs)
    marcadores.append((11, "BeepBox", "Logo BeepBox o captura de la composición musical"))
    diapositiva_dificultad_gif(prs)
    diapositiva_red(prs)
    marcadores.append((13, "Multijugador LAN", "Captura menú con IP o foto de 2 PCs jugando"))
    diapositiva_evolucion(prs, capturas)
    diapositiva_demo(prs, capturas)
    diapositiva_ia(prs)
    diapositiva_conclusiones(prs)
    diapositiva_preguntas(prs)

    prs.save(SALIDA)
    escribir_notas_marcadores(marcadores)
    print(f"Presentación: {SALIDA}")
    print(f"Diapositivas: {len(prs.slides)}")
    print(f"Marcadores: {NOTAS}")
    if capturas:
        print(f"Capturas evolución extraídas: {len(capturas)}")


if __name__ == "__main__":
    main()
